# -*- coding: utf-8 -*-
"""Export reveal.js deck slides to PNG + PPTX (full-bleed, no animations)."""
from pathlib import Path
import http.server
import socketserver
import threading

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Emu

DECK = Path(__file__).resolve().parent
SHOTS = DECK / "export-slides"
OUT_PPTX = DECK.parent / "defense-presentation-from-html.pptx"
OUT_PPTX_CN = DECK.parent / "智慧光棚-答辩演示-HTML.pptx"
PORT = 8765


def serve():
    handler = http.server.SimpleHTTPRequestHandler
    httpd = socketserver.TCPServer(("127.0.0.1", PORT), handler)
    httpd.allow_reuse_address = True
    t = threading.Thread(target=httpd.serve_forever, daemon=True)
    t.start()
    return httpd


def main():
    SHOTS.mkdir(exist_ok=True)
    old = Path.cwd()
    try:
        import os
        os.chdir(DECK)
        httpd = serve()
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={"width": 1280, "height": 720}, device_scale_factor=2)
            page.goto(f"http://127.0.0.1:{PORT}/index.html?export", wait_until="networkidle")
            page.wait_for_timeout(1500)
            # wait fonts/images
            page.evaluate("() => document.fonts.ready")
            page.wait_for_timeout(800)
            n = page.evaluate("() => Reveal.getTotalSlides()")
            print("SLIDES", n)
            paths = []
            for i in range(n):
                page.evaluate(f"() => Reveal.slide({i})")
                page.wait_for_timeout(500)
                path = SHOTS / f"slide-{i+1:02d}.png"
                page.locator(".reveal .slides").screenshot(path=str(path))
                # prefer full viewport for ppt
                page.screenshot(path=str(path), full_page=False)
                paths.append(path)
                print("OK", path.name, path.stat().st_size)
            browser.close()
        httpd.shutdown()
    finally:
        os.chdir(old)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for path in paths:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(OUT_PPTX))
    prs.save(str(OUT_PPTX_CN))
    print("WROTE", OUT_PPTX)
    print("WROTE", OUT_PPTX_CN)
    print("SIZE", OUT_PPTX.stat().st_size)


if __name__ == "__main__":
    main()
