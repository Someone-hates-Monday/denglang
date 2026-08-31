# -*- coding: utf-8 -*-
"""智慧光棚 答辩PPT 生成脚本 (python-pptx)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 主题色
DARK   = RGBColor(0x1E, 0x3A, 0x2E)   # 深绿
GREEN  = RGBColor(0x2F, 0x5D, 0x4A)   # 主绿
LIGHT  = RGBColor(0xEE, 0xF5, 0xEF)   # 浅绿底
GREY   = RGBColor(0x6B, 0x7D, 0x72)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x4C, 0x8A, 0x6B)

ASSETS = r"C:\Users\15374\Downloads\zard\denglang\docs\greenhouse\assets"
OUT    = r"C:\Users\15374\Downloads\zard\denglang\docs\greenhouse\智慧光棚-答辩PPT.pptx"

FONT = "微软雅黑"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def add_rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, lines, size=18, color=GREY, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [(lines, {})]
    for i, (txt, kw) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = kw.get("align", align)
        p.line_spacing = kw.get("line_spacing", line_spacing)
        run = p.add_run(); run.text = txt
        f = run.font
        f.name = FONT
        f.size = Pt(kw.get("size", size))
        f.bold = kw.get("bold", bold)
        f.color.rgb = kw.get("color", color)
    return tb

def title_bar(slide, num, title, subtitle=None):
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.14), GREEN)
    add_text(slide, Inches(0.6), Inches(0.35), Inches(9.5), Inches(0.7),
             [(title, {"size": 30, "bold": True, "color": DARK})])
    if subtitle:
        add_text(slide, Inches(0.62), Inches(1.0), Inches(12), Inches(0.4),
                 [(subtitle, {"size": 14, "color": GREY})])

def footer(slide, n):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(12), Inches(0.3),
             [("智慧光棚 · 设施农业光环境闭环管控平台    %d" % n,
               {"size": 10, "color": RGBColor(0xAA, 0xAA, 0xAA)})])

# ---------- 1 封面 ----------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, DARK)
add_rect(s, 0, Inches(4.9), prs.slide_width, Inches(0.06), ACCENT)
add_text(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.4),
         [("智慧光棚", {"size": 60, "bold": True, "color": WHITE})], align=PP_ALIGN.CENTER)
add_text(s, Inches(1.0), Inches(3.2), Inches(11.3), Inches(0.8),
         [("设施农业光环境闭环管控平台", {"size": 28, "color": RGBColor(0xCF,0xE0,0xD4)})],
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.0),
         [("测光 → 配方比对 → 补光/遮阳联动 → 农艺审批工单可审计",
           {"size": 18, "color": RGBColor(0xCF,0xE0,0xD4)})], align=PP_ALIGN.CENTER)
add_text(s, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.6),
         [("答辩人：XXX · 2026-08", {"size": 14, "color": RGBColor(0x9A,0xB5,0xA3)})],
         align=PP_ALIGN.CENTER)

# ---------- 2 背景与痛点 ----------
s = prs.slides.add_slide(BLANK)
title_bar(s, 2, "背景与痛点", "为什么需要一套光环境管控平台")
items = [
    ("重庆寡照", "年均日照时数全国偏低、冬春多雾，设施作物易长期欠光"),
    ("人工凭经验", "补光/遮阳靠人看天操作，动作滞后、难以复现、无法追责"),
    ("光配方缺位", "作物不同阶段需要不同 PPFD / DLI，缺乏量化目标带"),
    ("失控风险", "夏季强光暴晒、光抑制灼苗；大开度误操作缺审批保护"),
]
y = Inches(1.5)
for t, d in items:
    add_rect(s, Inches(0.8), y, Inches(0.22), Inches(0.5), ACCENT)
    add_text(s, Inches(1.2), y - Inches(0.03), Inches(11), Inches(0.55),
             [(t + "：", {"size": 20, "bold": True, "color": DARK})])
    add_text(s, Inches(1.2), y + Inches(0.42), Inches(11), Inches(0.4),
             [(d, {"size": 15, "color": GREY})])
    y += Inches(1.32)
footer(s, 2)

# ---------- 3 产品定位 ----------
s = prs.slides.add_slide(BLANK)
title_bar(s, 3, "产品定位", "面向设施栽培的光环境闭环管控平台")
cards = [
    ("主场景", "重庆近郊单跨拱棚\n16×7 m · 南向采光"),
    ("主作物", "铁皮石斛\n组培 / 栽培分阶段"),
    ("辅作物", "设施草莓 · 金线莲\n配方一键切换对比"),
    ("差异化", "空间光场 3D · 重庆日型仿真\n分床调光 · 审批≠执行"),
]
x = Inches(0.7); cw = Inches(2.85); gap = Inches(0.25)
for t, d in cards:
    add_rect(s, x, Inches(1.8), cw, Inches(2.9), LIGHT)
    add_rect(s, x, Inches(1.8), cw, Inches(0.7), GREEN)
    add_text(s, x, Inches(1.9), cw, Inches(0.5), [(t, {"size": 20, "bold": True, "color": WHITE})],
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(2.7), cw - Inches(0.3), Inches(1.8),
             [(d, {"size": 16, "color": DARK})], align=PP_ALIGN.CENTER)
    x += cw + gap
add_text(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.2),
         [("一句话：面向设施栽培的光环境闭环 —— 测光、比对、联动、审批、审计。",
           {"size": 18, "bold": True, "color": GREEN})])
footer(s, 3)

# ---------- 4 核心业务闭环 ----------
s = prs.slides.add_slide(BLANK)
title_bar(s, 4, "核心业务闭环", "测光 → 配方比对 → 补光/遮阳 → 工单审计")
steps = ["① 测光\nPAR 测点 + MQTT\n遥测 5s", "② 配方比对\n目标带 / 硬限\n区有效 PPFD", "③ 调控联动\n欠光补光\n过光遮阳优先", "④ 工单审计\n大开度≥80 审批\napprove≠execute"]
x = Inches(0.6); cw = Inches(2.6); gap = Inches(0.55)
for i, st in enumerate(steps):
    add_rect(s, x, Inches(1.7), cw, Inches(2.2), GREEN)
    add_text(s, x, Inches(1.95), cw, Inches(1.8), [(st, {"size": 18, "bold": True, "color": WHITE})],
             align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(s, x + cw + Inches(0.02), Inches(2.5), Inches(0.5), Inches(0.5),
                 [("→", {"size": 30, "bold": True, "color": ACCENT})], align=PP_ALIGN.CENTER)
    x += cw + gap
add_rect(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(1.9), LIGHT)
add_text(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(1.4),
         [("小步动作 AUTO 直接下发；", {"size": 16, "bold": True, "color": DARK}),
          ("大步动作（调光/遮阳 ≥ 80%）→ 生成 PENDING 工单 → 农艺师批准 → 种植员接单执行 → 控制日志与状态机全程可审计，避免误操作。",
           {"size": 16, "color": DARK})])
footer(s, 4)

# ---------- 5 作物光配方体系 ----------
s = prs.slides.add_slide(BLANK)
title_bar(s, 5, "作物光配方体系", "量化目标带 · 按需切换 · 解耦绑定")
cols = ["配方 / 阶段", "硬限下限", "目标带 PPFD", "硬限上限", "DLI 目标", "光周期"]
rows = [
    ["铁皮石斛（组培）", "50", "60–70", "90", "2.16–3.02", "12h"],
    ["铁皮石斛（栽培）", "70", "90–120", "140", "3.89–5.18", "12h"],
    ["设施草莓（冬春）", "150", "250–400", "550", "17–25", "12h"],
    ["台湾金线莲", "15", "25–35", "55", "1.26–1.76", "14h"],
]
tx = Inches(0.7); ty = Inches(1.7); cw = Inches(2.03)
add_rect(s, tx, ty, Inches(12.0), Inches(0.55), GREEN)
for ci, c in enumerate(cols):
    add_text(s, tx + ci*cw + Inches(0.05), ty + Inches(0.06), cw - Inches(0.1), Inches(0.4),
             [(c, {"size": 15, "bold": True, "color": WHITE})], align=PP_ALIGN.CENTER)
for ri, row in enumerate(rows):
    yy = ty + Inches(0.6 + ri*0.55)
    add_rect(s, tx, yy, Inches(12.0), Inches(0.52),
             LIGHT if ri % 2 == 0 else WHITE)
    for ci, cell in enumerate(row):
        add_text(s, tx + ci*cw + Inches(0.05), yy + Inches(0.06), cw - Inches(0.1), Inches(0.4),
                 [(cell, {"size": 14, "bold": ci in (2, 4), "color": DARK if ci != 2 else GREEN})],
                 align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.0),
         [("切换作物/阶段 = 改绑定 recipeId，无需改代码；DLI ≈ PPFD × 光周期(h) × 0.0036。",
           {"size": 15, "color": GREY})])
footer(s, 5)

# ---------- 6 空间光场与仿真 ----------
s = prs.slides.add_slide(BLANK)
title_bar(s, 6, "空间光场与重庆日型仿真", "3D 棚体 + 冠层热力 + 一天压缩 120s")
feats = [
    ("空间光场 3D", "Three.js 3D 棚架 + 灯位 + 冠层 PPFD 热力切片，非 2D 色块"),
    ("重庆日型仿真", "冬雾 / 夏正午两套日型，演示日压缩至约 2 分钟"),
    ("分床控光", "ZONE-A 石斛 / ZONE-B 金线莲分床独立目标带，互不干扰"),
    ("曲线与光谱", "全日光照/温湿度曲线，三色光谱通道，配方切换即时更新"),
]
y = Inches(1.6)
for t, d in feats:
    add_rect(s, Inches(0.8), y, Inches(0.22), Inches(0.5), ACCENT)
    add_text(s, Inches(1.2), y - Inches(0.03), Inches(11), Inches(0.5),
             [(t, {"size": 20, "bold": True, "color": DARK})])
    add_text(s, Inches(1.2), y + Inches(0.45), Inches(11), Inches(0.4),
             [(d, {"size": 15, "color": GREY})])
    y += Inches(1.3)
footer(s, 6)

# ---------- 7 角色分权与工单审计 ----------
s = prs.slides.add_slide(BLANK)
title_bar(s, 7, "角色分权与工单审计", "六角色 RBAC · 审批≠执行 · 全流程可追溯")
roles = [("场长", "策略启停 · 总览 · AUTO"), ("农艺师", "配方 CRUD · 工单审批"),
         ("种植员", "接单执行 · 现场回填"), ("设备运维", "设备档案 · 标定调试"),
         ("学员", "只读沙箱"), ("系统管理员", "账号与系统")]
x = Inches(0.7); cw = Inches(1.9); gap = Inches(0.12)
for i, (r, d) in enumerate(roles):
    add_rect(s, x, Inches(1.7), cw, Inches(1.6), LIGHT)
    add_rect(s, x, Inches(1.7), cw, Inches(0.55), GREEN)
    add_text(s, x, Inches(1.78), cw, Inches(0.4), [(r, {"size": 18, "bold": True, "color": WHITE})],
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), Inches(2.4), cw - Inches(0.2), Inches(0.8),
             [(d, {"size": 13, "color": DARK})], align=PP_ALIGN.CENTER)
    x += cw + gap
add_rect(s, Inches(0.7), Inches(3.7), Inches(11.9), Inches(2.6), LIGHT)
add_text(s, Inches(1.0), Inches(3.9), Inches(11.3), Inches(2.2),
         [("工单状态机：", {"size": 18, "bold": True, "color": DARK}),
          ("PENDING → APPROVED → EXECUTING → DONE ／ REJECTED", {"size": 17, "color": GREEN}),
          ("", {}),
          ("安全边界：农艺师 approve 仅审批；真正下发发生在种植员执行阶段（source=WORK_ORDER）。", {"size": 15, "color": DARK}),
          ("越权调用 403 · 控制日志不可物理删除 · 种植员无配方硬限写权限", {"size": 15, "color": DARK})])
footer(s, 7)

# ---------- 8 功能架构图 ----------
s = prs.slides.add_slide(BLANK)
title_bar(s, 8, "功能架构图", "四层：呈现 → 应用 → 数据 → 接入")
s.shapes.add_picture(ASSETS + r"\functional-architecture.png",
                     Inches(0.7), Inches(1.35), Inches(11.9), Inches(5.6))
footer(s, 8)

# ---------- 9 技术架构图 ----------
s = prs.slides.add_slide(BLANK)
title_bar(s, 9, "技术架构图", "Vue3 + Spring Boot + PostgreSQL(pgvector) + EMQX")
s.shapes.add_picture(ASSETS + r"\technical-architecture.png",
                     Inches(0.7), Inches(1.35), Inches(11.9), Inches(5.6))
footer(s, 9)

# ---------- 10 核心模块与实现 ----------
s = prs.slides.add_slide(BLANK)
title_bar(s, 10, "核心模块与实现", "后端 Java 21 / Spring Boot 3.5.9 · 前端 Vue3 / Vite 8")
mods = [
    ("后端 · gh 域", "配方引擎 · 控光规则 · 工单状态机 · 重庆日型仿真 · MQTT 上下行"),
    ("前端 · 冠层光场", "Three.js 3D · ECharts 曲线 · STOMP 实时 · RBAC 显隐"),
    ("数据 · PostgreSQL", "gh_zones / gh_recipes / gh_work_orders / 控制日志 + RAG 向量"),
    ("接入 · EMQX", "smart-greenhouse/{sn}/telemetry|status|command 契约"),
    ("农艺助手", "悬浮问答 · 27 条知识库 + RAG 检索 · 只读建议"),
]
y = Inches(1.6)
for t, d in mods:
    add_rect(s, Inches(0.8), y, Inches(2.2), Inches(0.75), GREEN)
    add_text(s, Inches(0.8), y + Inches(0.12), Inches(2.2), Inches(0.5),
             [(t, {"size": 16, "bold": True, "color": WHITE})], align=PP_ALIGN.CENTER)
    add_rect(s, Inches(3.1), y, Inches(9.6), Inches(0.75), LIGHT)
    add_text(s, Inches(3.3), y + Inches(0.1), Inches(9.2), Inches(0.55),
             [(d, {"size": 15, "color": DARK})], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.95)
footer(s, 10)

# ---------- 11 演示要点 ----------
s = prs.slides.add_slide(BLANK)
title_bar(s, 11, "3 分钟演示要点", "主叙事：铁皮石斛 · 光环境闭环")
demos = [
    "① 登录（admin / admin123），进入冠层光场，PPFD 曲线实时推进",
    "② 展示重庆日型仿真：冬雾寡照 → 自动补光；正午强光 → 遮阳优先",
    "③ 切换金线莲配方，目标带变为 25–35，无需改代码",
    "④ 制造大开度变更 → 生成 PENDING 工单 → 农艺批准 → 种植员执行 → 日志可查",
    "⑤ 角色切换：不同账号看到不同导航与按钮",
]
y = Inches(1.6)
for i, d in enumerate(demos):
    add_rect(s, Inches(0.8), y, Inches(0.5), Inches(0.5), GREEN)
    add_text(s, Inches(0.8), y + Inches(0.06), Inches(0.5), Inches(0.4),
             [(str(i+1), {"size": 18, "bold": True, "color": WHITE})], align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.5), y + Inches(0.02), Inches(10.8), Inches(0.6),
             [(d, {"size": 16, "color": DARK})], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.95)
footer(s, 11)

# ---------- 12 部署与验收 ----------
s = prs.slides.add_slide(BLANK)
title_bar(s, 12, "部署与验收", "从零到跑通 · 部署文档 DEPLOY.md")
steps = ["WSL2 + Docker\nPG :5433 / EMQX :1883", "git clone +\napplication-secret.yml", "init-db.ps1\n16 个 SQL 迁移",
         "后端编译启动\nJava 21 :8080", "前端 npm run build\nweb/dist :5173"]
x = Inches(0.5); cw = Inches(2.28); gap = Inches(0.14)
for i, st in enumerate(steps):
    add_rect(s, x, Inches(1.7), cw, Inches(1.5), LIGHT)
    add_text(s, x + Inches(0.1), Inches(1.85), cw - Inches(0.2), Inches(1.2),
             [(st, {"size": 14, "bold": True, "color": DARK})], align=PP_ALIGN.CENTER)
    x += cw + gap
add_rect(s, Inches(0.6), Inches(3.7), Inches(12.1), Inches(2.3), LIGHT)
add_text(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(2.0),
         [("验收清单（DEPLOY §8）：", {"size": 17, "bold": True, "color": DARK}),
          ("Docker 两容器健康 · gh_* 表就绪 · 登录 200 · 分区接口 2 区 · 冠层 PPFD 推进 · 工单闭环 · EMQX 控制台可登录",
           {"size": 15, "color": DARK}),
          ("", {}),
          ("关键排障：WSL 未装→Docker 500 · .env.local 带 BOM→mock 失效 · 端口占用 · 指令超时告警",
           {"size": 14, "color": GREY})])
footer(s, 12)

# ---------- 13 团队分工 ----------
s = prs.slides.add_slide(BLANK)
title_bar(s, 13, "团队分工（6 人）", "双轨协作 · RACI 对齐")
teams = [("组长 · 后端控光与仿真", "光场仿真 · 控光规则 · 工单门控 · MQTT"),
         ("前端 · 冠层与总览", "3D 光场 · 总览 · 控光交互 · 农艺助手"),
         ("数据库与数据", "gh_* 结构 · 配方种子 · 迁移 · 数据库文件"),
         ("角色权限与协同页", "六角色 RBAC · 导航显隐 · 申请/联系/报告"),
         ("需求 · 原型 · 架构图", "需求清单 · 思维导图 · 功能/技术架构图"),
         ("部署 · 测试 · 交付", "部署文档 · 发布包 · 联调回归 · PPT/视频")]
x = Inches(0.7); cw = Inches(3.9); gap = Inches(0.15); y = Inches(1.7)
for i, (t, d) in enumerate(teams):
    xx = x + (i % 2) * (cw + gap)
    yy = y + (i // 2) * Inches(1.35)
    add_rect(s, xx, yy, cw, Inches(1.2), LIGHT)
    add_rect(s, xx, yy, Inches(0.14), Inches(1.2), ACCENT)
    add_text(s, xx + Inches(0.3), yy + Inches(0.12), cw - Inches(0.4), Inches(0.45),
             [(t, {"size": 16, "bold": True, "color": DARK})])
    add_text(s, xx + Inches(0.3), yy + Inches(0.6), cw - Inches(0.4), Inches(0.5),
             [(d, {"size": 13, "color": GREY})])
footer(s, 13)

# ---------- 14 总结与展望 ----------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, DARK)
add_text(s, Inches(1.0), Inches(1.6), Inches(11.3), Inches(1.0),
         [("总结与展望", {"size": 40, "bold": True, "color": WHITE})], align=PP_ALIGN.CENTER)
lines = [
    "已完成：光环境闭环 · 空间光场 3D · 重庆日型仿真 · 角色分权工单审计 · 部署文档",
    "下一步：真机一路（RS-GZ / BH1750）· VPD/EC 多传感门控 · 冠层均匀度告警",
    "感谢聆听，欢迎提问",
]
y = Inches(3.0)
for ln in lines:
    add_text(s, Inches(1.2), y, Inches(10.9), Inches(0.7),
             [(ln, {"size": 20, "color": RGBColor(0xCF,0xE0,0xD4)})], align=PP_ALIGN.CENTER)
    y += Inches(1.0)

prs.save(OUT)
print("PPT saved:", OUT)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
