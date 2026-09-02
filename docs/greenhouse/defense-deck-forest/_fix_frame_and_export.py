# -*- coding: utf-8 -*-
from pathlib import Path
import re
import shutil

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches
import http.server
import socketserver
import threading
import os

ROOT = Path(__file__).resolve().parent
html_path = ROOT / "index.html"
html = html_path.read_text(encoding="utf-8")

# Add has-shot class to image slides
def add_has_shot(match: re.Match) -> str:
    open_tag = match.group(1)
    body = match.group(2)
    if "shot-stage" not in body:
        return match.group(0)
    if "has-shot" in open_tag:
        return match.group(0)
    open_tag = open_tag.replace('class="slide ', 'class="slide has-shot ', 1)
    return open_tag + body + "</section>"


html2, n = re.subn(
    r'(<section class="slide [^"]*"[^>]*>)([\s\S]*?)</section>',
    add_has_shot,
    html,
)
html_path.write_text(html2, encoding="utf-8")
print("patched has-shot replacements", n, "count", html2.count("has-shot"))

# Re-render diagrams
img = ROOT / "img"
diagram = (ROOT / "diagrams" / "forest-diagrams.html").as_uri()
with sync_playwright() as p:
    browser = p.chromium.launch(headless=True)
    page = browser.new_page(viewport={"width": 1920, "height": 1080})
    page.goto(diagram + "?page=loop", wait_until="networkidle")
    page.wait_for_timeout(400)
    page.screenshot(path=str(img / "14-flow-loop.png"))
    page.goto(diagram + "?page=arch", wait_until="networkidle")
    page.wait_for_timeout(400)
    page.screenshot(path=str(img / "15-arch-forest.png"))
    print("diagrams ok")
    browser.close()

# Export PPTX
SHOTS = ROOT / "export-slides"
SHOTS.mkdir(exist_ok=True)
OUT = ROOT.parent / "defense-presentation-forest-v2.pptx"
PORT = 8767


def serve():
    httpd = socketserver.TCPServer(("127.0.0.1", PORT), http.server.SimpleHTTPRequestHandler)
    httpd.allow_reuse_address = True
    threading.Thread(target=httpd.serve_forever, daemon=True).start()
    return httpd


old = Path.cwd()
os.chdir(ROOT)
httpd = serve()
try:
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": 1920, "height": 1080}, device_scale_factor=1)
        page.goto(f"http://127.0.0.1:{PORT}/index.html", wait_until="domcontentloaded")
        page.wait_for_timeout(1800)
        page.evaluate(
            """() => {
              document.body.classList.add('low-power');
              document.querySelectorAll('[data-anim]').forEach(el => {
                el.style.opacity = '1'; el.style.transform = 'none';
              });
              const hint = document.getElementById('hint');
              if (hint) hint.style.display = 'none';
              const nav = document.getElementById('nav');
              if (nav) nav.style.display = 'none';
            }"""
        )
        nslides = page.locator("section.slide").count()
        paths = []
        for i in range(nslides):
            page.evaluate("(i) => { if (typeof go === 'function') go(i, {force:true}); }", i)
            page.wait_for_timeout(450)
            path = SHOTS / f"slide-{i+1:02d}.png"
            page.screenshot(path=str(path), full_page=False)
            paths.append(path)
            print("OK", path.name)
        browser.close()
finally:
    httpd.shutdown()
    os.chdir(old)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
for path in paths:
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(str(path), 0, 0, width=prs.slide_width, height=prs.slide_height)
try:
    prs.save(str(OUT))
except PermissionError:
    OUT = ROOT.parent / "defense-presentation-forest-v3.pptx"
    prs.save(str(OUT))
print("WROTE", OUT, OUT.stat().st_size)

# update submit pack
pack4 = ROOT.parents[1] / "项目提交清单" / "4.ppt与答辩视频"
if pack4.exists():
    shutil.copy2(OUT, pack4 / "智慧光棚-答辩演示.pptx")
    print("updated submit pack")
